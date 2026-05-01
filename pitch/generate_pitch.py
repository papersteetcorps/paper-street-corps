"""Generate ForgeSelf pitch deck as a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand palette (from rebranding-to-forge branch) ────────────────────────
BG = RGBColor(0x0D, 0x0D, 0x14)
FG = RGBColor(0xF5, 0xF5, 0xF7)
SURFACE_800 = RGBColor(0x16, 0x16, 0x1F)
SURFACE_700 = RGBColor(0x2A, 0x2A, 0x3A)
SURFACE_500 = RGBColor(0x8E, 0x8E, 0x99)
SURFACE_400 = RGBColor(0x99, 0x99, 0xA8)
SURFACE_300 = RGBColor(0xBB, 0xBB, 0xC4)
ACCENT_AMBER = RGBColor(0xE8, 0x62, 0x2A)
ACCENT_PURPLE = RGBColor(0xA8, 0x55, 0xF7)
ACCENT_BLUE = RGBColor(0x7C, 0x5C, 0xFC)
ACCENT_TEAL = RGBColor(0x2D, 0xD4, 0xBF)
ACCENT_HIGHLIGHT = RGBColor(0xF0, 0xB4, 0x29)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Setup presentation (16:9 widescreen) ───────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height
BLANK_LAYOUT = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────
def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, x, y, w, h, text, *,
             font="Plus Jakarta Sans", size=14, bold=False,
             color=FG, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             letter_spacing=None, line_spacing=1.4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if letter_spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(letter_spacing))
    return tb


def add_eyebrow(slide, x, y, text, color=SURFACE_500):
    return add_text(slide, x, y, Inches(8), Inches(0.3), text.upper(),
                    font="Space Grotesk", size=11, bold=True, color=color,
                    letter_spacing=300)


def add_h1(slide, x, y, w, text, size=48, color=FG):
    return add_text(slide, x, y, w, Inches(2), text,
                    font="Space Grotesk", size=size, bold=True, color=color,
                    letter_spacing=-30, line_spacing=1.05)


def add_h2(slide, x, y, w, text, size=36, color=FG):
    return add_text(slide, x, y, w, Inches(1.5), text,
                    font="Space Grotesk", size=size, bold=True, color=color,
                    letter_spacing=-25, line_spacing=1.1)


def add_lead(slide, x, y, w, text, color=SURFACE_400):
    return add_text(slide, x, y, w, Inches(1.2), text,
                    font="Plus Jakarta Sans", size=16, color=color,
                    line_spacing=1.5)


def add_card(slide, x, y, w, h, *, fill=SURFACE_800, border=SURFACE_700,
             border_w=0.75, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(shape_type, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_w)
    if radius:
        card.adjustments[0] = 0.08
    card.shadow.inherit = False
    return card


def add_bar(slide, x, y, w, h, color):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def add_brand_corner(slide):
    add_text(slide, Inches(0.5), Inches(0.3), Inches(3), Inches(0.4),
             "Forge", font="Instrument Sans", size=14, bold=True, color=FG)
    add_text(slide, Inches(0.95), Inches(0.3), Inches(3), Inches(0.4),
             "Self", font="Instrument Sans", size=14, bold=True, color=ACCENT_AMBER)


def add_slide_num(slide, n, total):
    add_text(slide, Inches(11.5), Inches(0.3), Inches(1.5), Inches(0.3),
             f"{n:02d} / {total:02d}",
             font="Space Grotesk", size=10, color=SURFACE_500,
             align=PP_ALIGN.RIGHT, letter_spacing=200)


def add_divider(slide, x, y, w):
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Emu(9525))
    div.fill.solid()
    div.fill.fore_color.rgb = ACCENT_AMBER
    div.line.fill.background()
    div.shadow.inherit = False
    return div


def chrome(slide, n, total):
    add_brand_corner(slide)
    add_slide_num(slide, n, total)


TOTAL = 13


# ──────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ──────────────────────────────────────────────────────────────────────────
def slide_1():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)

    # Decorative accent rectangles (instead of aurora glow)
    glow1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(-2), Inches(8), Inches(8))
    glow1.fill.solid()
    glow1.fill.fore_color.rgb = RGBColor(0x3A, 0x1F, 0x4A)
    glow1.line.fill.background()
    glow1.shadow.inherit = False

    glow2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(3), Inches(8), Inches(7))
    glow2.fill.solid()
    glow2.fill.fore_color.rgb = RGBColor(0x4A, 0x25, 0x18)
    glow2.line.fill.background()
    glow2.shadow.inherit = False

    chrome(slide, 1, TOTAL)

    # FORGE
    add_text(slide, Inches(0), Inches(1.8), Inches(13.333), Inches(2),
             "FORGE", font="Instrument Sans", size=140, bold=True,
             color=ACCENT_AMBER, align=PP_ALIGN.CENTER, letter_spacing=-60,
             line_spacing=0.9)
    # SELF
    add_text(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(2),
             "SELF", font="Instrument Sans", size=140, bold=True,
             color=ACCENT_PURPLE, align=PP_ALIGN.CENTER, letter_spacing=-60,
             line_spacing=0.9)

    # divider
    add_bar(slide, Inches(6.166), Inches(5.45), Inches(1), Emu(15000), ACCENT_AMBER)

    # domain
    add_text(slide, Inches(0), Inches(5.65), Inches(13.333), Inches(0.5),
             "forgeself.xyz", font="Instrument Sans", size=18, bold=True,
             color=SURFACE_400, align=PP_ALIGN.CENTER)

    # tagline
    add_text(slide, Inches(0), Inches(6.2), Inches(13.333), Inches(0.6),
             "Figure yourself out. Then do something about it.",
             font="Plus Jakarta Sans", size=18, color=SURFACE_300,
             align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────
# Slide 2 — The Problem
# ──────────────────────────────────────────────────────────────────────────
def slide_2():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    chrome(slide, 2, TOTAL)

    add_eyebrow(slide, Inches(0.8), Inches(1.0), "The Problem")
    add_h1(slide, Inches(0.8), Inches(1.4), Inches(11.5),
           "Personality tests are broken.", size=52)
    add_lead(slide, Inches(0.8), Inches(2.7), Inches(10.5),
             "Most are either rigid multiple-choice or fluffy quizzes that tell you nothing useful. None of them help you actually change.")

    # Three problem cards
    points = [
        ("Multiple choice can't capture you",
         "Real personality is messy. Forcing it into A/B/C/D loses everything that matters."),
        ("Single frameworks are limited",
         "MBTI alone misses energy patterns. Enneagram alone misses cognitive style. You need multiple lenses."),
        ("No one tells you what to do next",
         '"You\'re an INTJ" — and? Where\'s the part where you actually grow?'),
    ]
    y = Inches(4.0)
    for title, body in points:
        card = add_card(slide, Inches(0.8), y, Inches(11.7), Inches(0.95))
        add_text(slide, Inches(1.05), y + Inches(0.13), Inches(0.4), Inches(0.4),
                 "→", font="Space Grotesk", size=20, bold=True, color=ACCENT_AMBER)
        add_text(slide, Inches(1.6), y + Inches(0.15), Inches(10.5), Inches(0.35),
                 title, font="Space Grotesk", size=15, bold=True, color=FG)
        add_text(slide, Inches(1.6), y + Inches(0.5), Inches(10.5), Inches(0.4),
                 body, font="Plus Jakarta Sans", size=12, color=SURFACE_400)
        y += Inches(1.05)


# ──────────────────────────────────────────────────────────────────────────
# Slide 3 — The Solution
# ──────────────────────────────────────────────────────────────────────────
def slide_3():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    chrome(slide, 3, TOTAL)

    add_eyebrow(slide, Inches(0.8), Inches(1.0), "The Solution")
    add_h1(slide, Inches(0.8), Inches(1.4), Inches(11.5),
           "Tell your story.\nWe read the patterns.", size=46)
    add_lead(slide, Inches(0.8), Inches(3.2), Inches(11.5),
             "ForgeSelf uses six psychological frameworks and AI to read across your life, find what's consistent, and tell you what to do about it.")

    # Three solution cards
    cards = [
        ("◈", "No multiple choice", "You write or speak about yourself. AI reads what's actually there."),
        ("◉", "Six frameworks", "Cross-validate your traits across Jungian, Socionics, Enneagram and more."),
        ("⚡", "Actually change", "Get concrete next steps, not vague horoscope advice."),
    ]
    cw, gap = Inches(3.85), Inches(0.25)
    start_x = Inches(0.8)
    y = Inches(4.6)
    for i, (icon, title, body) in enumerate(cards):
        x = start_x + (cw + gap) * i
        add_card(slide, x, y, cw, Inches(2.4))
        # icon box
        ib = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), y + Inches(0.3), Inches(0.55), Inches(0.55))
        ib.fill.solid()
        ib.fill.fore_color.rgb = RGBColor(0x33, 0x1A, 0x12)
        ib.line.fill.background()
        ib.adjustments[0] = 0.25
        ib.shadow.inherit = False
        add_text(slide, x + Inches(0.3), y + Inches(0.32), Inches(0.55), Inches(0.55),
                 icon, font="Space Grotesk", size=18, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.3), y + Inches(1.0), cw - Inches(0.6), Inches(0.5),
                 title, font="Space Grotesk", size=16, bold=True, color=FG)
        add_text(slide, x + Inches(0.3), y + Inches(1.45), cw - Inches(0.6), Inches(1),
                 body, font="Plus Jakarta Sans", size=12, color=SURFACE_400, line_spacing=1.5)


# ──────────────────────────────────────────────────────────────────────────
# Slide 4 — How It Works
# ──────────────────────────────────────────────────────────────────────────
def slide_4():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    chrome(slide, 4, TOTAL)

    add_eyebrow(slide, Inches(0.8), Inches(1.0), "How It Works")
    add_h2(slide, Inches(0.8), Inches(1.4), Inches(11.5),
           "Four steps to forge yourself.")

    steps = [
        ("01", "Tell your story", "Write or speak naturally. No checkboxes, no scales."),
        ("02", "AI reads patterns", "Six frameworks analyse what you wrote — across cognitive, emotional, and energy domains."),
        ("03", "See yourself clearly", "Get a unified profile that combines all six lenses into one honest picture."),
        ("04", "Forge ahead", "Concrete growth steps tailored to your patterns, plus a chat to dig deeper."),
    ]
    cw, gap = Inches(2.85), Inches(0.2)
    start_x = Inches(0.8)
    y = Inches(3.4)
    for i, (num, title, body) in enumerate(steps):
        x = start_x + (cw + gap) * i
        add_card(slide, x, y, cw, Inches(2.8))
        add_text(slide, x + Inches(0.3), y + Inches(0.3), cw - Inches(0.6), Inches(0.4),
                 num, font="Space Grotesk", size=12, bold=True, color=ACCENT_AMBER, letter_spacing=300)
        add_text(slide, x + Inches(0.3), y + Inches(0.85), cw - Inches(0.6), Inches(0.5),
                 title, font="Space Grotesk", size=15, bold=True, color=FG)
        add_text(slide, x + Inches(0.3), y + Inches(1.4), cw - Inches(0.6), Inches(1.3),
                 body, font="Plus Jakarta Sans", size=11.5, color=SURFACE_400, line_spacing=1.5)


# ──────────────────────────────────────────────────────────────────────────
# Slide 5 — Six Frameworks
# ──────────────────────────────────────────────────────────────────────────
def slide_5():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    chrome(slide, 5, TOTAL)

    add_eyebrow(slide, Inches(0.8), Inches(1.0), "What You Get")
    add_h2(slide, Inches(0.8), Inches(1.4), Inches(11.5),
           "Six frameworks. One you.")

    frameworks = [
        ("◈", "Jungian Type", "8 cognitive functions, 16 type stack, story-based.", ACCENT_BLUE),
        ("◉", "Socionics", "Information processing & interaction style.", ACCENT_AMBER),
        ("⚡", "Energy Profile", "Where your energy goes & burnout patterns.", ACCENT_PURPLE),
        ("⬡", "Temperaments", "Choleric, Melancholic, Phlegmatic, Sanguine.", ACCENT_PURPLE),
        ("⚖", "Moral Alignment", "Where you fall on order-chaos and good-evil.", ACCENT_TEAL),
        ("◎", "Enneagram", "Life-phase analysis across 9 types.", ACCENT_AMBER),
    ]
    cw, ch = Inches(3.85), Inches(1.65)
    gap = Inches(0.25)
    start_x = Inches(0.8)
    start_y = Inches(3.0)
    for i, (icon, title, body, color) in enumerate(frameworks):
        col = i % 3
        row = i // 3
        x = start_x + (cw + gap) * col
        y = start_y + (ch + gap) * row
        add_card(slide, x, y, cw, ch)
        # icon
        ib = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.3), y + Inches(0.25), Inches(0.5), Inches(0.5))
        ib.fill.solid()
        ib.fill.fore_color.rgb = SURFACE_700
        ib.line.fill.background()
        ib.adjustments[0] = 0.25
        ib.shadow.inherit = False
        add_text(slide, x + Inches(0.3), y + Inches(0.27), Inches(0.5), Inches(0.5),
                 icon, font="Space Grotesk", size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.3), y + Inches(0.85), cw - Inches(0.6), Inches(0.4),
                 title, font="Space Grotesk", size=15, bold=True, color=FG)
        add_text(slide, x + Inches(0.3), y + Inches(1.2), cw - Inches(0.6), Inches(0.45),
                 body, font="Plus Jakarta Sans", size=11, color=SURFACE_400, line_spacing=1.4)


# ──────────────────────────────────────────────────────────────────────────
# Slide 6 — The Forge Method
# ──────────────────────────────────────────────────────────────────────────
def slide_6():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    chrome(slide, 6, TOTAL)

    add_eyebrow(slide, Inches(0.8), Inches(1.0), "The Forge Method")
    add_h2(slide, Inches(0.8), Inches(1.4), Inches(11.5),
           "How to forge yourself.")
    add_lead(slide, Inches(0.8), Inches(2.4), Inches(11.5),
             "Knowing your type is step one. Forging yourself is the rest.")

    items = [
        ("Heat", "see yourself", "Honest input from six frameworks. The truths you've been avoiding surface in plain language."),
        ("Hammer", "name the pattern", "AI synthesis spots what runs through every framework — the core pattern beneath everything."),
        ("Shape", "act on it", 'Concrete growth path with specific actions. Not "try journaling" — actual moves for your type.'),
        ("Temper", "repeat & track", "Come back over time. Re-test, see how you've changed, sharpen the next iteration."),
    ]
    cw, ch = Inches(5.9), Inches(1.6)
    gap = Inches(0.25)
    start_x = Inches(0.8)
    start_y = Inches(3.4)
    for i, (verb, sub, body) in enumerate(items):
        col = i % 2
        row = i // 2
        x = start_x + (cw + gap) * col
        y = start_y + (ch + gap) * row
        add_card(slide, x, y, cw, ch)
        add_text(slide, x + Inches(0.35), y + Inches(0.2), cw - Inches(0.7), Inches(0.5),
                 f"{verb} — {sub}", font="Space Grotesk", size=16, bold=True, color=ACCENT_AMBER)
        add_text(slide, x + Inches(0.35), y + Inches(0.75), cw - Inches(0.7), Inches(0.85),
                 body, font="Plus Jakarta Sans", size=12, color=SURFACE_300, line_spacing=1.5)


# ──────────────────────────────────────────────────────────────────────────
# Slide 7 — Differentiation
# ──────────────────────────────────────────────────────────────────────────
def slide_7():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    chrome(slide, 7, TOTAL)

    add_eyebrow(slide, Inches(0.8), Inches(1.0), "Why ForgeSelf")
    add_h2(slide, Inches(0.8), Inches(1.4), Inches(11.5),
           "No one else does this.")

    headers = ["Feature", "Free quizzes", "16Personalities", "ForgeSelf"]
    rows = [
        ("Open-ended input", "—", "—", "✓"),
        ("Multiple frameworks", "—", "—", "6 frameworks"),
        ("AI synthesis", "—", "—", "✓"),
        ("Conversational chat", "—", "—", "✓"),
        ("Voice input", "—", "—", "✓"),
        ("Growth guidance", "—", "Generic", "Personalised"),
    ]

    table_x = Inches(0.8)
    table_y = Inches(2.7)
    col_widths = [Inches(3.2), Inches(2.7), Inches(2.7), Inches(3.1)]
    row_h = Inches(0.55)
    header_h = Inches(0.5)

    # Header row
    cx = table_x
    for i, h in enumerate(headers):
        add_text(slide, cx + Inches(0.15), table_y, col_widths[i], header_h,
                 h.upper(), font="Space Grotesk", size=10.5, bold=True,
                 color=SURFACE_500, letter_spacing=200)
        cx += col_widths[i]

    # Header divider
    add_bar(slide, table_x, table_y + Inches(0.4), sum(col_widths, Emu(0)), Emu(9525), SURFACE_700)

    # Rows
    ry = table_y + Inches(0.55)
    for ridx, row in enumerate(rows):
        # Row background card
        add_card(slide, table_x, ry, sum(col_widths, Emu(0)), row_h, fill=SURFACE_800, border_w=0)
        cx = table_x
        for i, cell in enumerate(row):
            color = FG
            if i == 3 and cell != "—":
                color = ACCENT_AMBER
            elif cell == "—":
                color = SURFACE_500
            add_text(slide, cx + Inches(0.15), ry + Inches(0.13), col_widths[i], row_h,
                     cell, font="Plus Jakarta Sans", size=12.5,
                     bold=(i == 3 and cell != "—"), color=color)
            cx += col_widths[i]
        ry += row_h + Inches(0.08)


# ──────────────────────────────────────────────────────────────────────────
# Slide 8 — Monetization
# ──────────────────────────────────────────────────────────────────────────
def slide_8():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    chrome(slide, 8, TOTAL)

    add_eyebrow(slide, Inches(0.8), Inches(1.0), "Monetization")
    add_h2(slide, Inches(0.8), Inches(1.4), Inches(11.5),
           "Freemium → premium.")
    add_lead(slide, Inches(0.8), Inches(2.4), Inches(11.5),
             "Free taste. Pay for depth.")

    items = [
        ("Hook", "One free framework. Strong enough to share, deep enough to want more."),
        ("Convert", "Unlock all 6 frameworks + AI synthesis + chat for $5/mo."),
        ("Retain", "Pro tier ($8/mo) adds voice, progress tracking, and personalised growth plans."),
    ]
    cw, gap = Inches(3.85), Inches(0.25)
    start_x = Inches(0.8)
    y = Inches(3.6)
    for i, (title, body) in enumerate(items):
        x = start_x + (cw + gap) * i
        add_card(slide, x, y, cw, Inches(2.5))
        add_text(slide, x + Inches(0.35), y + Inches(0.35), cw - Inches(0.7), Inches(0.5),
                 title, font="Space Grotesk", size=20, bold=True, color=ACCENT_AMBER)
        add_text(slide, x + Inches(0.35), y + Inches(1.0), cw - Inches(0.7), Inches(1.4),
                 body, font="Plus Jakarta Sans", size=13, color=SURFACE_300, line_spacing=1.55)


# ──────────────────────────────────────────────────────────────────────────
# Slide 9 — Pricing
# ──────────────────────────────────────────────────────────────────────────
def slide_9():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    chrome(slide, 9, TOTAL)

    add_eyebrow(slide, Inches(0.8), Inches(1.0), "Pricing")
    add_h2(slide, Inches(0.8), Inches(1.4), Inches(11.5),
           "Simple. Honest. Cheap.")

    plans = [
        ("Free", "$0", "/forever", [
            "1 framework (Jungian)", "Basic result page",
            "3 chat messages", "Email-only support"
        ], False),
        ("Forge", "$5", "/month", [
            "All 6 frameworks", "AI synthesis profile",
            "Unlimited chat", "Shareable result cards", "Re-test anytime"
        ], True),
        ("Forge Pro", "$8", "/month", [
            "Everything in Forge", "Voice input",
            "Progress tracking", "Personalised growth plan", "Priority support"
        ], False),
    ]

    cw, gap = Inches(3.85), Inches(0.25)
    start_x = Inches(0.8)
    y = Inches(2.8)
    for i, (name, price, suffix, feats, featured) in enumerate(plans):
        x = start_x + (cw + gap) * i
        # featured card raised slightly
        cy = y - Inches(0.15) if featured else y
        ch = Inches(4.3) if featured else Inches(4.0)
        fill = RGBColor(0x1F, 0x10, 0x0A) if featured else SURFACE_800
        border = ACCENT_AMBER if featured else SURFACE_700
        bw = 1.5 if featured else 0.75
        add_card(slide, x, cy, cw, ch, fill=fill, border=border, border_w=bw)

        if featured:
            # "Most Popular" badge
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + cw - Inches(1.6), cy - Inches(0.18), Inches(1.4), Inches(0.36))
            badge.fill.solid()
            badge.fill.fore_color.rgb = ACCENT_AMBER
            badge.line.fill.background()
            badge.adjustments[0] = 0.5
            badge.shadow.inherit = False
            add_text(slide, x + cw - Inches(1.6), cy - Inches(0.13), Inches(1.4), Inches(0.3),
                     "MOST POPULAR", font="Space Grotesk", size=9, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER, letter_spacing=200)

        # Plan name
        add_text(slide, x + Inches(0.4), cy + Inches(0.4), cw - Inches(0.8), Inches(0.4),
                 name, font="Space Grotesk", size=18, bold=True, color=FG)
        # Price
        add_text(slide, x + Inches(0.4), cy + Inches(0.95), cw - Inches(0.8), Inches(0.7),
                 price, font="Space Grotesk", size=42, bold=True, color=FG)
        add_text(slide, x + Inches(1.7), cy + Inches(1.45), cw - Inches(2), Inches(0.4),
                 suffix, font="Plus Jakarta Sans", size=13, color=SURFACE_500)

        # Features
        fy = cy + Inches(2.1)
        for feat in feats:
            add_text(slide, x + Inches(0.4), fy, Inches(0.3), Inches(0.3),
                     "→", font="Space Grotesk", size=12, bold=True, color=ACCENT_AMBER)
            add_text(slide, x + Inches(0.7), fy, cw - Inches(1.1), Inches(0.3),
                     feat, font="Plus Jakarta Sans", size=11.5, color=SURFACE_300)
            fy += Inches(0.34)


# ──────────────────────────────────────────────────────────────────────────
# Slide 10 — Revenue 100 users
# ──────────────────────────────────────────────────────────────────────────
def revenue_slide(slide, n, title, sub, rows, totals, stats):
    add_bg(slide)
    chrome(slide, n, TOTAL)
    add_eyebrow(slide, Inches(0.8), Inches(1.0), "Revenue Projection")
    add_h2(slide, Inches(0.8), Inches(1.4), Inches(11.5), title)
    add_lead(slide, Inches(0.8), Inches(2.4), Inches(11.5), sub)

    # Table
    headers = ["Plan", "Users", "Price", "Monthly", "Yearly"]
    col_widths = [Inches(3.0), Inches(2.0), Inches(2.0), Inches(2.4), Inches(2.4)]
    table_x = Inches(0.8)
    table_y = Inches(3.3)
    row_h = Inches(0.5)

    cx = table_x
    for i, h in enumerate(headers):
        add_text(slide, cx + Inches(0.15), table_y, col_widths[i], Inches(0.4),
                 h.upper(), font="Space Grotesk", size=10, bold=True,
                 color=SURFACE_500, letter_spacing=200)
        cx += col_widths[i]

    add_bar(slide, table_x, table_y + Inches(0.35), sum(col_widths, Emu(0)), Emu(9525), SURFACE_700)

    ry = table_y + Inches(0.5)
    for row in rows:
        add_card(slide, table_x, ry, sum(col_widths, Emu(0)), row_h, fill=SURFACE_800, border_w=0)
        cx = table_x
        for i, cell in enumerate(row):
            add_text(slide, cx + Inches(0.15), ry + Inches(0.1), col_widths[i], row_h,
                     cell, font="Plus Jakarta Sans", size=12.5, color=FG)
            cx += col_widths[i]
        ry += row_h + Inches(0.06)

    # Total row
    add_card(slide, table_x, ry, sum(col_widths, Emu(0)), row_h,
             fill=RGBColor(0x2D, 0x16, 0x0B), border=ACCENT_AMBER, border_w=1)
    cx = table_x
    for i, cell in enumerate(totals):
        add_text(slide, cx + Inches(0.15), ry + Inches(0.1), col_widths[i], row_h,
                 cell, font="Space Grotesk", size=13, bold=True, color=ACCENT_AMBER)
        cx += col_widths[i]

    # Stats
    sw, sgap = Inches(3.85), Inches(0.25)
    sx = Inches(0.8)
    sy = Inches(6.05)
    for i, (val, lbl) in enumerate(stats):
        x = sx + (sw + sgap) * i
        add_card(slide, x, sy, sw, Inches(1.15))
        add_text(slide, x, sy + Inches(0.2), sw, Inches(0.6),
                 val, font="Space Grotesk", size=26, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)
        add_text(slide, x, sy + Inches(0.78), sw, Inches(0.3),
                 lbl, font="Plus Jakarta Sans", size=11.5, color=SURFACE_400, align=PP_ALIGN.CENTER)


def slide_10():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    revenue_slide(slide, 10, "100 paying users.",
                  "Assumed split: 70% Forge ($5) / 30% Forge Pro ($8).",
                  rows=[
                      ("Forge", "70", "$5", "$350", "$4,200"),
                      ("Forge Pro", "30", "$8", "$240", "$2,880"),
                  ],
                  totals=("Total", "100", "—", "$590", "$7,080"),
                  stats=[("$590", "MRR"), ("$7,080", "ARR"), ("$5.90", "Avg ARPU")])


def slide_11():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    revenue_slide(slide, 11, "500 paying users.",
                  "Same split. The numbers scale.",
                  rows=[
                      ("Forge", "350", "$5", "$1,750", "$21,000"),
                      ("Forge Pro", "150", "$8", "$1,200", "$14,400"),
                  ],
                  totals=("Total", "500", "—", "$2,950", "$35,400"),
                  stats=[("$2,950", "MRR"), ("$35,400", "ARR"), ("5x", "vs 100 users")])


# ──────────────────────────────────────────────────────────────────────────
# Slide 12 — Team & Equity
# ──────────────────────────────────────────────────────────────────────────
def slide_12():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    chrome(slide, 12, TOTAL)

    add_eyebrow(slide, Inches(0.8), Inches(1.0), "Team & Economics")
    add_h2(slide, Inches(0.8), Inches(1.4), Inches(11.5),
           "Who builds this. What it costs.")
    add_lead(slide, Inches(0.8), Inches(2.4), Inches(11.5),
             "Three founders. Two minimum salaries. One mission.")

    members = [
        ("Founder & Lead", "Zoe", "40%", "No salary", "Equity-only", True),
        ("Co-founder & Engineering", "Gourav", "30%", "Monthly salary", "$600 / mo", False),
        ("Co-founder & Engineering", "Pratik", "30%", "Monthly salary", "$600 / mo", False),
    ]
    cw, gap = Inches(3.85), Inches(0.25)
    start_x = Inches(0.8)
    y = Inches(3.2)
    for i, (role, name, eq, sal_label, sal_val, lead) in enumerate(members):
        x = start_x + (cw + gap) * i
        fill = RGBColor(0x1F, 0x10, 0x0A) if lead else SURFACE_800
        border = ACCENT_AMBER if lead else SURFACE_700
        bw = 1.5 if lead else 0.75
        add_card(slide, x, y, cw, Inches(2.6), fill=fill, border=border, border_w=bw)

        add_text(slide, x, y + Inches(0.3), cw, Inches(0.3),
                 role.upper(), font="Space Grotesk", size=10, bold=True,
                 color=SURFACE_500, align=PP_ALIGN.CENTER, letter_spacing=200)
        add_text(slide, x, y + Inches(0.65), cw, Inches(0.5),
                 name, font="Space Grotesk", size=22, bold=True, color=FG, align=PP_ALIGN.CENTER)
        add_text(slide, x, y + Inches(1.2), cw, Inches(0.7),
                 eq, font="Space Grotesk", size=44, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)
        add_text(slide, x, y + Inches(1.85), cw, Inches(0.3),
                 "EQUITY", font="Space Grotesk", size=10, color=SURFACE_500, align=PP_ALIGN.CENTER, letter_spacing=200)
        # divider
        add_bar(slide, x + Inches(0.5), y + Inches(2.15), cw - Inches(1), Emu(9525), SURFACE_700)
        add_text(slide, x, y + Inches(2.22), cw, Inches(0.3),
                 sal_label, font="Plus Jakarta Sans", size=11, color=SURFACE_400, align=PP_ALIGN.CENTER)
        add_text(slide, x, y + Inches(2.42), cw, Inches(0.3),
                 sal_val, font="Space Grotesk", size=14, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)

    # Burn stats
    burn = [
        ("$1,200", "Monthly burn (salaries)"),
        ("~204", "Users to break even"),
        ("$1,750", "Net profit at 500 users / mo"),
    ]
    bw_, bgap = Inches(3.85), Inches(0.25)
    bx = Inches(0.8)
    by = Inches(6.1)
    for i, (val, lbl) in enumerate(burn):
        x = bx + (bw_ + bgap) * i
        add_card(slide, x, by, bw_, Inches(1.05))
        add_text(slide, x, by + Inches(0.18), bw_, Inches(0.5),
                 val, font="Space Grotesk", size=22, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)
        add_text(slide, x, by + Inches(0.7), bw_, Inches(0.3),
                 lbl, font="Plus Jakarta Sans", size=11, color=SURFACE_400, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────
# Slide 13 — Roadmap + CTA
# ──────────────────────────────────────────────────────────────────────────
def slide_13():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(slide)
    chrome(slide, 13, TOTAL)

    add_eyebrow(slide, Inches(0.8), Inches(1.0), "What's Next")
    add_h2(slide, Inches(0.8), Inches(1.4), Inches(11.5), "The roadmap.")

    cols = [
        ("Now (MVP)", ["6 frameworks live", "AI synthesis", "Result chat", "Browser voice input"]),
        ("Q3 2026", ["Shareable result cards", "Native mobile UX", "Stripe billing live", "Re-test history"]),
        ("Q4 2026", ["Progress tracking", "Personalised growth plans", "Premium voice (AssemblyAI)", "iOS / Android app"]),
        ("2027+", ["Coach marketplace", "Couples & team profiles", "Affiliate program", "API for researchers"]),
    ]
    cw, gap = Inches(2.95), Inches(0.18)
    start_x = Inches(0.8)
    y = Inches(2.5)
    for i, (h, items) in enumerate(cols):
        x = start_x + (cw + gap) * i
        add_card(slide, x, y, cw, Inches(2.6))
        add_text(slide, x + Inches(0.25), y + Inches(0.25), cw - Inches(0.5), Inches(0.4),
                 h.upper(), font="Space Grotesk", size=11, bold=True,
                 color=ACCENT_AMBER, letter_spacing=200)
        ly = y + Inches(0.75)
        for item in items:
            add_text(slide, x + Inches(0.25), ly, cw - Inches(0.5), Inches(0.4),
                     item, font="Plus Jakarta Sans", size=11.5, color=SURFACE_300)
            ly += Inches(0.42)

    # CTA box
    cx = Inches(0.8)
    cy = Inches(5.5)
    cwb = Inches(11.7)
    chb = Inches(1.7)
    cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cwb, chb)
    cta.fill.solid()
    cta.fill.fore_color.rgb = RGBColor(0x1F, 0x10, 0x0A)
    cta.line.color.rgb = ACCENT_AMBER
    cta.line.width = Pt(1.5)
    cta.adjustments[0] = 0.1
    cta.shadow.inherit = False

    add_text(slide, cx, cy + Inches(0.25), cwb, Inches(0.7),
             "FORGESELF", font="Instrument Sans", size=42, bold=True,
             color=ACCENT_AMBER, align=PP_ALIGN.CENTER, letter_spacing=-30)
    add_text(slide, cx, cy + Inches(0.95), cwb, Inches(0.4),
             "forgeself.xyz", font="Instrument Sans", size=16, bold=True,
             color=SURFACE_300, align=PP_ALIGN.CENTER)
    add_text(slide, cx, cy + Inches(1.3), cwb, Inches(0.35),
             "Figure yourself out. Then do something about it.",
             font="Plus Jakarta Sans", size=12, color=SURFACE_400, align=PP_ALIGN.CENTER)


# ── Build all slides ───────────────────────────────────────────────────────
slide_1()
slide_2()
slide_3()
slide_4()
slide_5()
slide_6()
slide_7()
slide_8()
slide_9()
slide_10()
slide_11()
slide_12()
slide_13()

out = "pitch/forgeself-pitch.pptx"
prs.save(out)
print(f"Generated {out}")
